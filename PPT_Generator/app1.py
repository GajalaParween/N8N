from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE

prs = Presentation()

# Slide Layouts:
# 0: Title Slide
# 1: Title and Content
# 3: Title Only
# Using Title and Content for most slides for consistency.

# Slide 1: Introduction to Air Quality Index (AQI)
slide_layout_1 = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Understanding the Air Quality Index (AQI)"
tf = body.text_frame
tf.text = "What is AQI? A simple numerical scale and color-coded system."
p = tf.add_paragraph()
p.text = "Purpose: To communicate daily air quality levels and associated health risks to the public."
p = tf.add_paragraph()
p.text = "Consolidates data from multiple pollutants into a single, easy-to-understand value."
p = tf.add_paragraph()
p.text = "Helps individuals make informed decisions to protect their health."
p = tf.add_paragraph()
p.text = "India's System: Based on 8 key pollutants for comprehensive assessment."

# Slide 2: How AQI is Measured
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Monitoring and Calculation of AQI"
tf = body.text_frame
tf.text = "Key Pollutants: PM2.5, PM10, Ozone (O3), Nitrogen Dioxide (NO2), Carbon Monoxide (CO), Sulfur Dioxide (SO2), Ammonia (NH3), Lead (Pb)."
p = tf.add_paragraph()
p.text = "Monitoring Stations: Automated continuous monitoring stations collect real-time data."
p = tf.add_paragraph()
p.text = "Sub-indices: A separate sub-index is calculated for each pollutant based on its concentration and impact."
p = tf.add_paragraph()
p.text = "Overall AQI: The highest sub-index value among all monitored pollutants becomes the city's AQI."
p = tf.add_paragraph()
p.text = "Data Processing: Raw data is converted into an AQI value and health descriptor."
p = tf.add_paragraph()
p.text = "(Visual suggestion: Infographic showing monitoring station components and data flow)."

# Slide 3: AQI Categories and Health Implications
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "AQI Categories and Associated Health Risks"
tf = body.text_frame
tf.text = "Good (0-50, Green): Minimal impact; ideal air quality."
p = tf.add_paragraph()
p.text = "Satisfactory (51-100, Light Green): Minor breathing discomfort to sensitive people possible."
p = tf.add_paragraph()
p.text = "Moderately Polluted (101-200, Yellow): Breathing discomfort to people with lung/heart disease, children, and older adults."
p = tf.add_paragraph()
p.text = "Poor (201-300, Orange): Breathing discomfort to most people on prolonged exposure."
p = tf.add_paragraph()
p.text = "Very Poor (301-400, Red): Respiratory illness on prolonged exposure; significant impact on vulnerable groups."
p = tf.add_paragraph()
p.text = "Severe (401-500+, Dark Red): Affects healthy people; serious health impacts on those with existing diseases."

# Slide 4: Major Causes of High AQI in India
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Sources of Air Pollution Driving AQI in India"
tf = body.text_frame
tf.text = "Vehicular Emissions: Rapid increase in private vehicles, old vehicle fleet, traffic congestion, poor fuel quality."
p = tf.add_paragraph()
p.text = "Industrial Pollution: Emissions from power plants, manufacturing units, brick kilns, and small-scale industries."
p = tf.add_paragraph()
p.text = "Construction Activities: Dust from building sites, infrastructure projects, and demolition."
p = tf.add_paragraph()
p.text = "Crop Residue Burning (Stubble Burning): Seasonal practice in agricultural regions, especially in North India."
p = tf.add_paragraph()
p.text = "Biomass Burning: Domestic use of wood, cow dung, and agricultural waste for cooking and heating."
p = tf.add_paragraph()
p.text = "Geographical & Meteorological Factors: Temperature inversions, low wind speed, and landlocked regions trapping pollutants."

# Slide 5: Analysis of AQI in Indian Cities
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Trends and Patterns of AQI Across Indian Cities"
tf = body.text_frame
tf.text = "Urbanization Impact: Rapid growth in cities leads to increased population density and pollution sources."
p = tf.add_paragraph()
p.text = "Regional Disparities: North Indian cities often experience higher pollution levels, particularly during winter."
p = tf.add_paragraph()
p.text = "Seasonal Variations: AQI typically worsens in winter due to calm winds and temperature inversions."
p = tf.add_paragraph()
p.text = "Case Study - Delhi: Consistently ranks among the world's most polluted cities due to a combination of factors."
p = tf.add_paragraph()
p.text = "Other Affected Cities: Kanpur, Ghaziabad, Lucknow, Patna, and others frequently record 'Very Poor' or 'Severe' AQI."
p = tf.add_paragraph()
p.text = "(Visual suggestion: Map highlighting cities with consistently high AQI across India)."

# Slide 6: Strategies to Reduce AQI
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Actionable Measures for Improving Air Quality"
tf = body.text_frame
tf.text = "Promote Public Transportation: Expand and improve bus, metro, and rail networks; encourage electric vehicles (EVs)."
p = tf.add_paragraph()
p.text = "Strict Emission Norms: Enforce advanced emission standards for vehicles (e.g., BS-VI) and industries."
p = tf.add_paragraph()
p.text = "Green Infrastructure: Large-scale afforestation, development of green belts, and urban forests."
p = tf.add_paragraph()
p.text = "Sustainable Waste Management: Implement proper waste segregation, composting, and ban open burning."
p = tf.add_paragraph()
p.text = "Renewable Energy Transition: Shift from coal-fired power plants to solar, wind, and hydropower."
p = tf.add_paragraph()
p.text = "Public Awareness & Participation: Educate citizens on pollution impacts and promote behavioral changes."

# Slide 7: Focus Areas for High AQI Cities
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Targeted Interventions for Heavily Polluted Cities"
tf = body.text_frame
tf.text = "Emergency Response Plans: Implement Graded Response Action Plan (GRAP) based on AQI levels (e.g., in Delhi-NCR)."
p = tf.add_paragraph()
p.text = "Pollution Hotspot Identification: Pinpoint major sources and deploy targeted, localized interventions."
p = tf.add_paragraph()
p.text = "Dust Control Measures: Regular road sweeping, water sprinkling, anti-smog guns, and strict construction site norms."
p = tf.add_paragraph()
p.text = "Industrial Upgrades: Mandate cleaner technologies, emissions control systems, and relocation of polluting industries."
p = tf.add_paragraph()
p.text = "Traffic Management: Congestion reduction, promotion of non-motorized transport, and vehicle restriction schemes."
p = tf.add_paragraph()
p.text = "(Visual suggestion: Infographic detailing various stages and actions of GRAP)."

# Slide 8: Focus Areas for Low AQI Cities
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Sustaining Good Air Quality in Less Polluted Cities"
tf = body.text_frame
tf.text = "Proactive Planning: Integrate air quality considerations into all urban development and industrial planning."
p = tf.add_paragraph()
p.text = "Green Development: Prioritize green buildings, sustainable transportation infrastructure, and open spaces."
p = tf.add_paragraph()
p.text = "Continuous Monitoring: Maintain and expand AQI monitoring networks to detect early changes or emerging issues."
p = tf.add_paragraph()
p.text = "Public Engagement: Encourage community participation in local environmental protection and green initiatives."
p = tf.add_paragraph()
p.text = "Learning from Others: Analyze challenges faced by high AQI cities to implement preventative measures."
p = tf.add_paragraph()
p.text = "(Visual suggestion: Images of green cities or sustainable urban planning projects)."

# Slide 9: Visuals and Graphical Analysis of AQI
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Visualizing Air Quality Trends and Impacts"
tf = body.text_frame
tf.text = "AQI Trend Line Graphs: Showcase changes in a city's AQI over seasons, months, or years."
p = tf.add_paragraph()
p.text = "Comparative Bar Charts: Illustrate AQI differences across major Indian cities (e.g., Delhi vs. Bengaluru vs. Chennai)."
p = tf.add_paragraph()
p.text = "Pollutant Contribution Pie Charts: Break down the percentage contribution of various pollutants to overall AQI."
p = tf.add_paragraph()
p.text = "Color-Coded AQI Maps: Provide an intuitive, real-time overview of air quality across regions or states."
p = tf.add_paragraph()
p.text = "Infographics on Health Impacts: Visually represent the effects of different AQI levels on human health and specific organs."
p = tf.add_paragraph()
p.text = "(Actual charts and images would be embedded here in a live presentation from data sources)."

# Slide 10: Conclusion on AQI
slide = prs.slides.add_slide(slide_layout_1)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Towards a Healthier Future: Concluding Thoughts on AQI"
tf = body.text_frame
tf.text = "AQI serves as a critical indicator for understanding and addressing the complex issue of air pollution."
p = tf.add_paragraph()
p.text = "India faces significant air quality challenges, particularly in rapidly urbanizing and industrialized areas."
p = tf.add_paragraph()
p.text = "A multi-faceted approach involving robust policies, technological advancements, and public participation is crucial."
p = tf.add_paragraph()
p.text = "Prioritizing sustainable development and green initiatives is key to long-term air quality improvement."
p = tf.add_paragraph()
p.text = "Collective responsibility and sustained action are essential for ensuring a cleaner, healthier India for all."

# Save the presentation
prs.save("generated_presentation.pptx")